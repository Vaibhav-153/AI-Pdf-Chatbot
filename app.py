# app.py - Modified for Streamlit Cloud Deployment

import os
import io
import base64
import mammoth
from pptx import Presentation
import streamlit as st
import google.generativeai as genai

# --- Backend Imports ---
from backend.ingestion import process_uploaded_files, get_text_chunks
from backend.embedding import create_hybrid_retriever
from backend.chat import get_rag_response

# --- CONFIGURATION ---
st.set_page_config(
    page_title="ManthanAI",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- SESSION STATE INIT ---
def initialize_session_state():
    if "api_keys_loaded" not in st.session_state:
        # Load keys from Streamlit secrets, which is the required method for deployment
        st.session_state.google_api_key = st.secrets.get("GEMINI_API_KEY")
        st.session_state.cohere_api_key = st.secrets.get("COHERE_API_KEY")
        st.session_state.api_keys_loaded = True
        if st.session_state.google_api_key:
            genai.configure(api_key=st.session_state.google_api_key)

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = [{"role": "ai", "content": "Hello! I am ManthanAI. Ask me anything or upload documents for analysis."}]
    if "retriever" not in st.session_state:
        st.session_state.retriever = None
    if "search_mode" not in st.session_state:
        st.session_state.search_mode = "Hybrid"
    if "uploaded_files_data" not in st.session_state:
        st.session_state.uploaded_files_data = {}
    if "selected_document" not in st.session_state:
        st.session_state.selected_document = None

# --- HELPER & FILE HANDLING FUNCTIONS ---
def display_pdf(file_bytes):
    """Encodes and displays a PDF file in an iframe."""
    base64_pdf = base64.b64encode(file_bytes).decode('utf-8')
    return f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="700" type="application/pdf"></iframe>'

def display_docx(file_bytes):
    """Extracts and displays text from a DOCX file."""
    try:
        # Use an in-memory buffer
        buffer = io.BytesIO(file_bytes)
        result = mammoth.convert_to_html(buffer)
        return result.value
    except Exception as e:
        return f"<p>Could not read DOCX file. Error: {e}</p>"

def display_pptx(file_bytes):
    """Extracts and displays text from a PPTX file."""
    try:
        buffer = io.BytesIO(file_bytes)
        prs = Presentation(buffer)
        full_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = [f"<h3>--- Slide {i+1} ---</h3>"]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(f"<p>{shape.text}</p>")
            full_text.append("".join(slide_text))
        return "".join(full_text)
    except Exception as e:
        return f"<p>Could not read PPTX file. Error: {e}</p>"


def handle_file_upload():
    if st.session_state.get('file_uploader_key'):
        with st.spinner("Processing uploaded documents..."):
            try:
                uploaded_files = st.session_state.file_uploader_key
                st.session_state.uploaded_files_data = {file.name: file.getvalue() for file in uploaded_files}
                
                raw_docs = process_uploaded_files(uploaded_files)
                text_chunks = get_text_chunks(raw_docs)
                if not text_chunks:
                    st.error("No text extracted. Please check your files.")
                    return
                st.session_state.retriever = create_hybrid_retriever(
                    text_chunks, st.session_state.google_api_key, st.session_state.cohere_api_key
                )
                st.toast(f"✅ Processed {len(st.session_state.uploaded_files_data)} documents successfully!")
            except Exception as e:
                st.error(f"Processing error: {e}")
                st.session_state.retriever = None
                st.session_state.uploaded_files_data = {}

# --- UI COMPONENTS ---
def build_sidebar():
    with st.sidebar:
        st.title("ManthanAI")
        st.caption("AI-powered document insights.")

        st.header("Step 1: Upload Files")
        st.file_uploader(
            "Upload PDFs, DOCX, or PPTX", type=["pdf", "docx", "pptx"],
            accept_multiple_files=True, key='file_uploader_key', on_change=handle_file_upload
        )

        if st.session_state.uploaded_files_data:
            st.header("Step 2: View Document")
            doc_list = ["--- Select a document ---"] + list(st.session_state.uploaded_files_data.keys())
            st.selectbox(
                "Choose a file to display:",
                options=doc_list,
                key="selected_document"
            )

        st.header("Step 3: Choose Mode")
        st.radio(
            "AI Mode:", ("Hybrid", "PDF-Only"), key="search_mode", horizontal=True
        )
        
        st.header("Step 4: Reset")
        if st.button("🔄 Reset Session", use_container_width=True):
            for key in list(st.session_state.keys()):
                if key not in ['api_keys_loaded', 'google_api_key', 'cohere_api_key']:
                    del st.session_state[key]
            initialize_session_state()
            st.rerun()

def display_chat_interface(tab):
    """Displays the chat history and the input box within a given tab."""
    with tab:
        st.header("🤖 ManthanAI - Your AI Document Analyst")
        st.caption("Ask anything from your PDF or just chat with AI.")
        
        with st.expander("⚠️ Gemini API Limit Info", expanded=False):
            st.info(
                "You're using the Free Tier of Gemini API, which allows only **50 requests/day**.\n"
                "If this limit is exceeded, ManthanAI may not respond until the next day.\n"
                "Sorry for the inconvenience caused."
                
             )
        
        
        for msg in st.session_state.chat_history:
            with st.chat_message(msg["role"], avatar="🤖" if msg["role"] == "ai" else "🧑"):
                st.markdown(msg["content"])
                
        return st.chat_input("Ask a question...")

def display_document_viewer(tab):
    """Displays the selected document in the viewer tab."""
    with tab:
        st.header("📄 Document Viewer")
        selected_doc = st.session_state.selected_document
        if selected_doc and selected_doc != "--- Select a document ---":
            file_bytes = st.session_state.uploaded_files_data.get(selected_doc)
            
            # Logic to display the correct viewer based on file type
            if selected_doc.lower().endswith(".pdf"):
                viewer_html = display_pdf(file_bytes)
                st.markdown(viewer_html, unsafe_allow_html=True)
            elif selected_doc.lower().endswith(".docx"):
                viewer_html = display_docx(file_bytes)
                st.markdown(viewer_html, unsafe_allow_html=True)
            elif selected_doc.lower().endswith(".pptx"):
                viewer_html = display_pptx(file_bytes)
                st.markdown(viewer_html, unsafe_allow_html=True)

        else:
            st.info("Select a document from the dropdown in the sidebar to view it here.")

# --- MAIN APPLICATION ---
def main():
    initialize_session_state()
    if not st.session_state.google_api_key or not st.session_state.cohere_api_key:
        st.error("Missing API keys. Please add GEMINI_API_KEY and COHERE_API_KEY to your Streamlit secrets.")
        st.stop()

    build_sidebar()

    chat_tab, viewer_tab = st.tabs(["💬 Chat", "📄 Document Viewer"])
    user_input = display_chat_interface(chat_tab)
    display_document_viewer(viewer_tab)
    
    greetings = ['hi', 'hello', 'hey']
    if user_input:
        st.session_state.chat_history.append({"role": "user", "content": user_input})
        
        if user_input.lower().strip() in greetings:
            st.session_state.chat_history.append({"role": "ai", "content": "Hello! How can I help you today?"})
        elif st.session_state.retriever:
            with st.spinner("Analyzing documents..."):
                response = get_rag_response(
                    user_input, st.session_state.retriever,
                    st.session_state.google_api_key, st.session_state.search_mode
                )
                st.session_state.chat_history.append({
                    "role": "ai", "content": response["output_text"],
                    "sources": response.get("source_documents", [])
                })
        else:
            with st.spinner("Thinking..."):
                try:
                    model = genai.GenerativeModel('gemini-pro')
                    response = model.generate_content(user_input)
                    st.session_state.chat_history.append({"role": "ai", "content": response.text})
                except Exception as e:
                    st.error(f"Could not get a response: {e}")
        
        st.rerun()

if __name__ == "__main__":
    main()