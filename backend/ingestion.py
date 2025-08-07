# pdf_chatbot/backend/ingestion.py

"""
Module for robust data ingestion from various document formats.

This upgraded module handles the initial step of the RAG pipeline: extracting 
raw text and metadata from uploaded files. It now supports PDF, DOCX, and PPTX formats.

Key Improvements:
1.  Uses PyMuPDF (fitz) for superior and faster PDF text extraction.
2.  Adds support for DOCX (Microsoft Word) and PPTX (Microsoft PowerPoint).
3.  A central dispatcher function `process_uploaded_files` routes files to the correct parser.
4.  Maintains metadata (source filename, page/slide number) for accurate citations.
"""

import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from typing import List, BinaryIO, Tuple
from langchain_core.documents import Document

# --- Text Chunking ---

def get_text_chunks(documents: List[Document]) -> List[Document]:
    """
    Splits a list of Documents into smaller, more manageable chunks.

    This is a key step for vector embedding, as it allows the model to find
    more specific and relevant pieces of information during a similarity search.
    The metadata from the parent document is automatically preserved in each chunk.

    Args:
        documents: A list of Document objects (from the parsing functions).

    Returns:
        A list of smaller Document chunks, ready for embedding.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        # The size of each text chunk. 1000 characters is a balanced choice.
        chunk_size=1000,
        # How many characters to overlap between chunks. This helps maintain
        # context so that a topic isn't split awkwardly between two chunks.
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_documents(documents)
    return chunks

# --- Document Parsers ---

def _parse_pdf(file_stream: BinaryIO, filename: str) -> List[Document]:
    """Extracts text from a PDF file stream using PyMuPDF."""
    documents = []
    try:
        with fitz.open(stream=file_stream, filetype="pdf") as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text:  # Ensure the page has text content
                    documents.append(Document(
                        page_content=text,
                        metadata={'source': filename, 'page': page_num + 1}
                    ))
    except Exception as e:
        print(f"Error parsing PDF '{filename}': {e}")
    return documents

def _parse_docx(file_stream: BinaryIO, filename: str) -> List[Document]:
    """Extracts text from a DOCX file stream."""
    documents = []
    try:
        doc = docx.Document(file_stream)
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text])
        if full_text:
            # For DOCX, we treat the whole file as one document initially
            documents.append(Document(
                page_content=full_text,
                metadata={'source': filename}
            ))
    except Exception as e:
        print(f"Error parsing DOCX '{filename}': {e}")
    return documents

def _parse_pptx(file_stream: BinaryIO, filename: str) -> List[Document]:
    """Extracts text from a PPTX file stream."""
    documents = []
    try:
        presentation = Presentation(file_stream)
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                documents.append(Document(
                    page_content=slide_text,
                    metadata={'source': filename, 'slide': slide_num + 1}
                ))
    except Exception as e:
        print(f"Error parsing PPTX '{filename}': {e}")
    return documents

# --- Main Ingestion Dispatcher ---

def process_uploaded_files(uploaded_files: List[BinaryIO]) -> List[Document]:
    """
    Processes a list of uploaded files, dispatching them to the correct parser.

    Args:
        uploaded_files: A list of uploaded file objects from Streamlit,
                        each having a 'name' and 'getvalue()' attribute.

    Returns:
        A list of LangChain Document objects, one for each page/slide/document.
    """
    all_documents = []
    if not uploaded_files:
        return all_documents

    for file_stream in uploaded_files:
        if file_stream is None:
            continue

        filename = file_stream.name
        file_extension = os.path.splitext(filename)[1].lower()
        
        # Reset stream position to the beginning before reading
        file_stream.seek(0)
        
        # Route to the appropriate parser based on file extension
        if file_extension == '.pdf':
            all_documents.extend(_parse_pdf(file_stream, filename))
        elif file_extension == '.docx':
            all_documents.extend(_parse_docx(file_stream, filename))
        elif file_extension == '.pptx':
            all_documents.extend(_parse_pptx(file_stream, filename))
        else:
            print(f"Warning: Unsupported file type '{file_extension}' for file '{filename}'. Skipping.")
            continue
            
    return all_documents